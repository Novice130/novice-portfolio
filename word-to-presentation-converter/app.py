import io
import re
import os
import uuid
import shutil
import subprocess
import tempfile
from typing import Optional
from flask import Flask, request, render_template, send_file, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import PyPDF2

app = Flask(__name__)

# Pattern to detect question numbers like "1.", "1)", "Q1", "Q.1", "(1)", etc.
QUESTION_NUM_PATTERN = re.compile(r'^\s*(?:Q\.?\s*)?(\d{1,3})\s*[.):]')
# Pattern for option labels (A., a), A), A:, etc.)
OPTION_LABEL_PATTERN = re.compile(r'^\s*([A-Da-d])\s*[.):]\s*')
# Inline options: e.g. "A) ... B) ... C) ... D) ..."
INLINE_OPTIONS_PATTERN = re.compile(r'([A-Da-d])\s*[.):]\s*', re.IGNORECASE)


_BULLET_PREFIX_RE = r"(?:[\u2022\u2023\u25E6\u2043\u2219\u00B7\u25CF\-\*])?"


_ANSWER_MARK_RE = re.compile(r"\bAnswer\b\s*:?", re.IGNORECASE)


def _strip_answer_artifacts(text: str) -> tuple[str, bool]:
    """Remove trailing answer-key artifacts.

    Returns: (cleaned_text, saw_answer_marker)
    """
    if not text:
        return "", False
    m = _ANSWER_MARK_RE.search(text)
    if not m:
        return text.strip(), False
    # Keep text before the answer marker; drop everything after.
    cleaned = text[: m.start()].rstrip(" .,:;\t")
    return cleaned.strip(), True


def _split_docx_paragraph_indices(doc: Document) -> list[list[int]]:
    """Split a DOCX into paragraph-index blocks using the same heuristic as text splitting."""
    lines = [_normalize_extracted_text(p.text.rstrip()) for p in doc.paragraphs]
    blocks: list[list[int]] = []
    buf: list[int] = []
    for idx, line in enumerate(lines):
        stripped = (line or "").strip()
        if stripped == "":
            if buf:
                blocks.append(buf)
                buf = []
            continue
        if QUESTION_NUM_PATTERN.match(stripped):
            if buf:
                blocks.append(buf)
                buf = []
        buf.append(idx)
    if buf:
        blocks.append(buf)
    return blocks


def _create_pptx_by_copy_paste_from_docx(docx_path: str, font_family: str, font_size: int) -> Optional[io.BytesIO]:
    """Windows-only: Use Word+PowerPoint COM automation to copy/paste question blocks.

    This preserves Word equation/fraction formatting far better than rebuilding text.
    Returns a BytesIO PPTX on success, or None if COM isn't available.
    """
    if os.name != 'nt':
        return None
    try:
        import win32com.client  # type: ignore
    except Exception:
        return None

    # Build paragraph blocks using python-docx so block boundaries match our normal logic.
    doc = Document(docx_path)
    para_blocks = _split_docx_paragraph_indices(doc)
    if not para_blocks:
        return None

    # Determine optional title block.
    first_block_text = "\n".join(
        [doc.paragraphs[i].text.strip() for i in para_blocks[0] if doc.paragraphs[i].text.strip()]
    ).strip()
    has_title = _is_probable_title_block(first_block_text)

    # Prepare output temp file
    fd, out_pptx = tempfile.mkstemp(suffix='.pptx')
    os.close(fd)

    word = None
    ppt = None
    wdoc = None
    pres = None
    try:
        word = win32com.client.Dispatch('Word.Application')
        word.Visible = False
        wdoc = word.Documents.Open(os.path.abspath(docx_path), ReadOnly=True)

        ppt = win32com.client.Dispatch('PowerPoint.Application')
        # Paste operations are more reliable with the app visible.
        try:
            ppt.Visible = True
        except Exception:
            pass

        pres = ppt.Presentations.Add()

        # Layout constants
        ppLayoutBlank = 12
        left = 40
        top = 40
        width = 880
        height = 480

        def add_slide_with_pasted_range(rng):
            slide = pres.Slides.Add(pres.Slides.Count + 1, ppLayoutBlank)
            rng.Copy()
            shp_range = slide.Shapes.Paste()
            shp = shp_range[1]
            try:
                shp.Left = left
                shp.Top = top
                shp.Width = width
                shp.Height = height
            except Exception:
                pass
            # Apply font choices; leave equation objects intact if PowerPoint keeps them separate.
            try:
                if hasattr(shp, 'TextFrame') and shp.TextFrame.HasText:
                    shp.TextFrame.TextRange.Font.Name = font_family
                    shp.TextFrame.TextRange.Font.Size = font_size
            except Exception:
                pass

        # Title slide (as plain text) if present.
        start_block_idx = 0
        if has_title and first_block_text:
            slide = pres.Slides.Add(pres.Slides.Count + 1, ppLayoutBlank)
            tx = slide.Shapes.AddTextbox(1, left, top, width, height)
            tx.TextFrame.TextRange.Text = first_block_text
            try:
                tx.TextFrame.TextRange.Font.Name = font_family
                tx.TextFrame.TextRange.Font.Size = min(44, font_size + 6)
            except Exception:
                pass
            start_block_idx = 1

        # Question slides via copy/paste from Word.
        for block in para_blocks[start_block_idx:]:
            if not block:
                continue

            # Remove any answer key section: stop at the first paragraph that starts with "Answer".
            filtered = []
            for i in block:
                t = (doc.paragraphs[i].text or '').strip()
                if re.match(r'^\s*Answer\b', t, re.I):
                    break
                filtered.append(i)
            if not filtered:
                continue

            start_i = filtered[0] + 1  # COM paragraphs are 1-based
            end_i = filtered[-1] + 1
            rng = wdoc.Range(wdoc.Paragraphs(start_i).Range.Start, wdoc.Paragraphs(end_i).Range.End)
            add_slide_with_pasted_range(rng)

        pres.SaveAs(os.path.abspath(out_pptx))

        with open(out_pptx, 'rb') as f:
            bio = io.BytesIO(f.read())
        bio.seek(0)
        return bio

    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if ppt is not None:
                ppt.Quit()
        except Exception:
            pass
        try:
            if wdoc is not None:
                wdoc.Close(False)
        except Exception:
            pass
        try:
            if word is not None:
                word.Quit()
        except Exception:
            pass
        try:
            os.remove(out_pptx)
        except Exception:
            pass


def _normalize_extracted_text(text: str) -> str:
    """Normalize common PDF text-extraction artifacts.

    Goals:
    - Preserve mixed fractions that often come out as spaced digits: "1 1 2" -> "1 1/2"
    - Preserve simple fractions that often come out as spaced digits: "34 35" -> "34/35"
    - Keep changes conservative by applying fraction reconstruction only on math-like lines.
    """
    if not text:
        return text

    # Normalize non-breaking/thin spaces that frequently appear in PDF extraction.
    text = (
        text.replace("\u00A0", " ")
        .replace("\u2009", " ")
        .replace("\u202F", " ")
        .replace("\u2007", " ")
    )

    # Normalize slash-like characters.
    text = (
        text.replace("⁄", "/")
        .replace("∕", "/")
        .replace("／", "/")
    )

    def looks_mathy(line: str) -> bool:
        if not line or not line.strip():
            return False
        lower = line.lower()
        if "what is" in lower or "answer" in lower:
            return True
        # option line (allow bullet prefix): "• A) ..."
        if re.match(rf"^\s*{_BULLET_PREFIX_RE}\s*[A-Da-d]\s*[.):]", line):
            return True
        # math operators or an explicit slash already
        return bool(re.search(r"[+\-−×*/=()\[\]]|/", line))

    mixed_spaced = re.compile(r"(?<!\d)(\d{1,3})\s+(\d{1,3})\s+(\d{1,3})(?!\d)")
    # Note: exclude cases where the second number is immediately followed by '/',
    # which would indicate we've already normalized to something like "1 1/2".
    frac_spaced = re.compile(r"(?<!\d)(\d{1,3})\s+(\d{1,3})(?![\d/])")

    def mixed_repl(m: re.Match) -> str:
        whole_s, numer_s, denom_s = m.group(1), m.group(2), m.group(3)
        try:
            whole = int(whole_s)
            numer = int(numer_s)
            denom = int(denom_s)
        except ValueError:
            return m.group(0)
        if denom <= 0:
            return m.group(0)
        # Keep it conservative: typical classroom denominators are not huge.
        if denom > 200 or numer > 200:
            return m.group(0)
        return f"{whole} {numer}/{denom}"

    def frac_repl(m: re.Match) -> str:
        numer_s, denom_s = m.group(1), m.group(2)
        try:
            numer = int(numer_s)
            denom = int(denom_s)
        except ValueError:
            return m.group(0)
        if denom <= 0:
            return m.group(0)
        if denom > 200 or numer > 999:
            return m.group(0)
        return f"{numer}/{denom}"

    out_lines = []
    for line in text.splitlines():
        # Collapse excessive whitespace but keep newlines.
        line = re.sub(r"[\t\r\f\v]+", " ", line)
        line = re.sub(r" {2,}", " ", line).strip()

        if looks_mathy(line):
            # First reconstruct mixed fractions (3-number runs), then simple fractions (2-number runs).
            line = mixed_spaced.sub(mixed_repl, line)
            line = frac_spaced.sub(frac_repl, line)

        out_lines.append(line)

    return "\n".join(out_lines).strip()


def _is_probable_title_block(block: str) -> bool:
    """Return True if the first block is likely a document title.

    Important: do NOT classify numbered questions like "1." as a title.
    """
    if not block:
        return False
    first_line = next((l.strip() for l in block.splitlines() if l.strip()), "")
    if not first_line:
        return False
    # If it looks like a question number (e.g., "1.", "Q1)"), it's not a title.
    if QUESTION_NUM_PATTERN.match(first_line):
        return False
    # If it looks like an option-only block (e.g., "A) 624"), it's not a title.
    if OPTION_LABEL_PATTERN.match(first_line):
        return False
    # Keep the existing heuristic: short first block is likely a title.
    return len(block.split()) <= 30


def split_into_blocks(doc):
    """Split document paragraphs into question blocks, handling both blank-line and question-number delimiters."""
    paragraphs = [_normalize_extracted_text(p.text.rstrip()) for p in doc.paragraphs]
    return _split_lines_into_blocks(paragraphs)


def _split_lines_into_blocks(lines):
    """Given a list of text lines, split into question blocks."""
    blocks = []
    buf = []
    for line in lines:
        stripped = line.strip()
        if stripped == "":
            if buf:
                blocks.append("\n".join(buf).strip())
                buf = []
            continue
        # Check if this line starts a new question (has question number at start)
        if QUESTION_NUM_PATTERN.match(stripped):
            if buf:
                blocks.append("\n".join(buf).strip())
                buf = []
        buf.append(line)
    if buf:
        blocks.append("\n".join(buf).strip())
    return blocks


def split_text_blocks(text):
    """Split raw text (e.g. from PDF) into question blocks."""
    text = _normalize_extracted_text(text)
    lines = text.splitlines()
    return _split_lines_into_blocks(lines)


def _parse_inline_options(text):
    """If text contains inline options like 'A) x B) y ...', extract them."""
    # Find all option labels
    matches = list(INLINE_OPTIONS_PATTERN.finditer(text))
    if len(matches) < 2:
        return None  # not enough labels
    options = []
    for i, m in enumerate(matches):
        label = m.group(1).upper()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        raw = text[start:end].strip().rstrip(' ,;.')
        cleaned, saw_answer = _strip_answer_artifacts(raw)
        if cleaned:
            options.append(f"{label}) {cleaned}")
        if saw_answer:
            break
    return options if len(options) >= 2 else None


def classify_block(block):
    """Classify a block as MCQ, fill-in, or short answer."""
    lines = [l.rstrip() for l in block.splitlines() if l.strip()]
    if not lines:
        return {'type': 'short', 'question': ''}

    # Check for inline options in entire block text
    full_text = ' '.join(lines)
    inline_opts = _parse_inline_options(full_text)
    if inline_opts:
        # Extract question text before first option
        first_match = INLINE_OPTIONS_PATTERN.search(full_text)
        q_text = full_text[:first_match.start()].strip() if first_match else full_text
        q_text, _ = _strip_answer_artifacts(q_text)
        return {'type': 'mcq', 'question': q_text, 'options': inline_opts}

    # Match options on separate lines like: "• A) ...", "A.", "A:", etc.
    option_pattern = re.compile(rf'^\s*{_BULLET_PREFIX_RE}\s*([A-Da-d])\s*(?:[.):])+\s*(.*)$')
    has_options = any(option_pattern.match(l) for l in lines)
    has_blank = any('_' in l for l in lines)

    if has_options:
        q_lines = []
        options = []
        in_answer_section = False
        for l in lines:
            if re.match(r'^\s*Answer\b', l, re.I):
                break
            m = option_pattern.match(l)
            if m:
                label = m.group(1).upper()
                raw = _normalize_extracted_text(m.group(2).strip())
                content, saw_answer = _strip_answer_artifacts(raw)
                if content:
                    options.append(f"{label}) {content}")
                if saw_answer:
                    in_answer_section = True
                    break
            else:
                if not options:
                    q_lines.append(_normalize_extracted_text(l.strip()))
            if in_answer_section:
                break
        question = ' '.join(q_lines).strip()
        return {'type': 'mcq', 'question': question or lines[0], 'options': options}
    if has_blank:
        return {'type': 'fill', 'question': ' '.join(lines)}
    return {'type': 'short', 'question': ' '.join(lines)}


def parse_docx(stream):
    doc = Document(stream)
    blocks = split_into_blocks(doc)
    title = None
    q_blocks = []
    if not blocks:
        return None, []
    # Heuristic: first block is title if short-ish (<=30 words) AND not a question.
    if _is_probable_title_block(blocks[0]):
        title = blocks[0]
        remaining = blocks[1:]
    else:
        remaining = blocks

    # Merge small following blocks that look like options into preceding question block
    merged = []
    opt_label_pattern = re.compile(r'^\s*[A-Za-z]\s*(?:[.\)])')
    i = 0
    while i < len(remaining):
        cur = remaining[i]
        j = i + 1
        # append subsequent blocks that start with option labels or are short option-like lines
        while j < len(remaining) and opt_label_pattern.search(remaining[j].splitlines()[0] if remaining[j].splitlines() else ''):
            cur = cur + '\n' + remaining[j]
            j += 1
        merged.append(cur)
        i = j
    q_blocks = merged

    questions = []
    for b in q_blocks:
        cls = classify_block(b)
        questions.append(cls)

    return title, questions


def parse_pdf(stream):
    reader = PyPDF2.PdfReader(stream)
    text = []
    for page in reader.pages:
        try:
            text.append(page.extract_text() or '')
        except Exception:
            text.append('')
    full = '\n\n'.join(text)
    blocks = split_text_blocks(full)
    title = None
    q_blocks = blocks
    if blocks and _is_probable_title_block(blocks[0]):
        title = blocks[0]
        q_blocks = blocks[1:]

    # Merge option-like following pieces into same question block (same logic as docx)
    opt_label_pattern = re.compile(r'^\s*[A-Za-z]\s*(?:[.\)])')
    merged = []
    i = 0
    while i < len(q_blocks):
        cur = q_blocks[i]
        j = i + 1
        while j < len(q_blocks) and opt_label_pattern.search(q_blocks[j].splitlines()[0] if q_blocks[j].splitlines() else ''):
            cur = cur + '\n' + q_blocks[j]
            j += 1
        merged.append(cur)
        i = j

    questions = [classify_block(b) for b in merged]
    return title, questions


def find_soffice():
    # Try to find LibreOffice soffice executable
    candidates = [
        shutil.which('soffice'),
        shutil.which('soffice.exe'),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return None


def convert_doc_to_docx(input_path):
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError('LibreOffice soffice not found. Install LibreOffice or add soffice to PATH')
    outdir = os.path.dirname(input_path)
    # run soffice headless to convert
    cmd = [soffice, '--headless', '--convert-to', 'docx', '--outdir', outdir, input_path]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"soffice conversion failed: {proc.stderr}")
    base = os.path.splitext(os.path.basename(input_path))[0]
    converted = os.path.join(outdir, base + '.docx')
    if not os.path.exists(converted):
        raise RuntimeError('Converted file not found')
    return converted


def save_upload_temp(file_storage, suffix=''):
    fd, path = tempfile.mkstemp(suffix=suffix)
    os.close(fd)
    file_storage.save(path)
    return path


def apply_font_to_paragraph(p, text, font_size_pt, font_family, spacing_mult, is_question=False):
    """Apply font consistently by clearing paragraph and adding a fresh run with explicit font settings."""
    # Clear existing content
    p.clear()
    # Add a run with explicit font settings
    run = p.add_run()
    run.text = text
    # Set font properties directly on the run (most reliable method)
    run.font.size = Pt(font_size_pt)
    run.font.name = font_family
    run.font.color.rgb = RGBColor(0, 0, 0)
    run.font.bold = False
    run.font.italic = False
    # Also set on paragraph level as fallback
    p.font.size = Pt(font_size_pt)
    p.font.name = font_family
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.bold = False
    p.font.italic = False
    # Spacing
    if is_question:
        p.space_after = Pt(int(8 * spacing_mult))
        p.space_before = Pt(int(4 * spacing_mult))
    else:
        p.space_before = Pt(int(4 * spacing_mult))
        p.space_after = Pt(int(4 * spacing_mult))


def verify_and_fix_pptx(prs, font_family, question_font_pt, option_font_pt, answer_font_pt):
    """
    Scan all slides and fix any font inconsistencies.
    Ensures all text uses the selected font and size.
    """
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Preserve special formatting (fractions built with super/subscript)
                    # by not forcing those runs to the paragraph target size.
                    is_scripted = bool(getattr(run.font, 'superscript', False) or getattr(run.font, 'subscript', False))

                    # Determine which font size this run should have
                    text_lower = run.text.lower()
                    
                    # Detect what type of content this is
                    if any(x in text_lower for x in ['answer:', 'answer ', 'answer_']):
                        target_size = Pt(answer_font_pt)
                    elif any(x in text_lower for x in ['a)', 'b)', 'c)', 'd)', 'a.', 'b.', 'c.', 'd.']):
                        target_size = Pt(option_font_pt)
                    else:
                        target_size = Pt(question_font_pt)
                    
                    # Force consistent font properties
                    if (not is_scripted) and run.font.size != target_size:
                        run.font.size = target_size
                    if run.font.name != font_family:
                        run.font.name = font_family
                    if run.font.color.rgb != RGBColor(0, 0, 0):
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    if run.font.bold:
                        run.font.bold = False
                    if run.font.italic:
                        run.font.italic = False


_FRACTION_TOKEN_RE = re.compile(
    r"(?P<mixed>(?P<whole>\d{1,4})\s+(?P<num>\d{1,4})\s*/\s*(?P<den>\d{1,4}))|(?P<simple>(?P<num2>\d{1,4})\s*/\s*(?P<den2>\d{1,4}))"
)


def _add_text_with_fraction_formatting(paragraph, text, font_family, font_pt):
    """Add text to a paragraph, formatting a/b as numerator superscript + fraction slash + denominator subscript.

    This approximates Word/PPT fraction appearance better than plain '1/2'.
    """
    if text is None:
        text = ''
    text = str(text)

    def add_plain(t: str, size_pt: int):
        if not t:
            return
        r = paragraph.add_run()
        r.text = t
        r.font.size = Pt(size_pt)
        r.font.name = font_family
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.bold = False
        r.font.italic = False

    def add_num(t: str, size_pt: int):
        if not t:
            return
        r = paragraph.add_run()
        r.text = t
        r.font.size = Pt(size_pt)
        r.font.name = font_family
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.bold = False
        r.font.italic = False
        r.font.superscript = True

    def add_den(t: str, size_pt: int):
        if not t:
            return
        r = paragraph.add_run()
        r.text = t
        r.font.size = Pt(size_pt)
        r.font.name = font_family
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.bold = False
        r.font.italic = False
        r.font.subscript = True

    # fraction pieces are typically smaller
    frac_pt = max(10, int(round(font_pt * 0.65)))

    pos = 0
    for m in _FRACTION_TOKEN_RE.finditer(text):
        start, end = m.span()
        add_plain(text[pos:start], font_pt)

        if m.group('mixed'):
            whole = m.group('whole')
            num = m.group('num')
            den = m.group('den')
            add_plain(whole + ' ', font_pt)
            add_num(num, frac_pt)
            add_plain('⁄', frac_pt)
            add_den(den, frac_pt)
        else:
            num = m.group('num2')
            den = m.group('den2')
            add_num(num, frac_pt)
            add_plain('⁄', frac_pt)
            add_den(den, frac_pt)

        pos = end

    add_plain(text[pos:], font_pt)


def create_pptx(title, questions, font_family='Calibri', font_size=36, line_spacing=1.25):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    
    # Convert parameters - USE CONSISTENT SIZES ACROSS ALL SLIDES
    base_font_size = int(font_size)
    spacing_mult = float(line_spacing)
    
    # Fixed font sizes based on user selection (consistent across all slides)
    question_font_pt = base_font_size
    option_font_pt = max(16, int(base_font_size * 0.72))
    answer_font_pt = max(18, int(base_font_size * 0.6))
    
    # Title slide
    if title:
        slide = prs.slides.add_slide(blank_layout)
        left = top = Inches(1)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(2)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        p = tf.paragraphs[0]
        # Use add_run for guaranteed font application
        p.clear()
        run = p.add_run()
        run.text = "\n\n" + title + "\n\n"
        run.font.size = Pt(min(40, base_font_size + 4))
        run.font.name = font_family
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.bold = False
        run.font.italic = False
        p.alignment = PP_ALIGN.CENTER

    # Layout params
    # base margins
    left = Inches(0.6)
    right = Inches(0.6)
    top_margin = Inches(0.4)
    # user requested small pixel offsets: 20px right, 80px down (moved further down)
    px_to_in = 1.0 / 96.0
    extra_right = Inches(20 * px_to_in)
    extra_down = Inches(80 * px_to_in)  # increased from 30 to 80
    bottom_margin = Inches(0.6)
    usable_width_in = (prs.slide_width - left - right) / Inches(1)
    usable_height_in = (prs.slide_height - top_margin - bottom_margin) / Inches(1)

    def estimate_height_inches(text, font_pt, width_in, avg_chars_per_inch=8, line_mult=1.2):
        if not text: 
            return 0.0
        chars = len(text)
        # average characters per inch heuristic — lower value -> more conservative wrapping
        chars_per_line = max(20, int(width_in * avg_chars_per_inch))
        lines = (chars + chars_per_line - 1) // chars_per_line
        height_in = lines * (font_pt * line_mult) / 72.0
        return height_in

    def split_text_to_fit(text, font_pt, width_in, max_height_in):
        # naive split by characters using estimate function — returns list of chunks that fit
        if not text:
            return []
        words = text.split()
        chunks = []
        cur = []
        for w in words:
            cur.append(w)
            h = estimate_height_inches(' '.join(cur), font_pt, width_in)
            if h > max_height_in:
                # pop last word and finalize chunk
                cur.pop()
                if cur:
                    chunks.append(' '.join(cur))
                cur = [w]
        if cur:
            chunks.append(' '.join(cur))
        return chunks

    # Question slides - use CONSISTENT font sizes
    for q in questions:
        slide = prs.slides.add_slide(blank_layout)

        q_text = _normalize_extracted_text(q.get('question', ''))
        opts = [_normalize_extracted_text(o) for o in (q.get('options', []) if q.get('type') == 'mcq' else [])]

        # detect and remove any explicit answer line(s) present in the question text
        existing_answer = None
        if q_text:
            q_lines = [ln for ln in q_text.splitlines() if ln.strip()]
            tail_lines = []
            while q_lines and (re.search(r'_{3,}', q_lines[-1]) or re.match(r'^\s*Answer[:\s]', q_lines[-1], re.I)):
                tail_lines.insert(0, q_lines.pop())
            if tail_lines:
                existing_answer = '\n'.join(tail_lines).strip()
            q_text = ' '.join(q_lines).strip()

        # apply pixel offsets
        place_left = left + extra_right
        place_top = top_margin + extra_down

        # Calculate if content fits, if not split across slides
        combined = q_text
        if q['type'] == 'mcq' and opts:
            combined = combined + '\n' + '\n'.join(opts)
        elif q['type'] in ('fill', 'short'):
            combined = combined + '\n' + (existing_answer or 'Answer:')

        est_h = estimate_height_inches(combined, question_font_pt, usable_width_in)
        
        # If content fits in one slide, render it; otherwise split
        if est_h <= usable_height_in:
            # place everything in one textbox with reduced font
            question_height = Inches(max(0.6, est_h))
            txq = slide.shapes.add_textbox(place_left, place_top, prs.slide_width - left - right - extra_right, question_height)
            tfq = txq.text_frame
            # ensure wrapping and small internal margins so text wraps inside box
            try:
                tfq.word_wrap = True
            except Exception:
                pass
            try:
                tfq.margin_left = Pt(6)
                tfq.margin_right = Pt(6)
                tfq.margin_top = Pt(6)
                tfq.margin_bottom = Pt(6)
            except Exception:
                pass
            tfq.clear()
            # set question paragraph with CONSISTENT font using add_run
            p_q = tfq.paragraphs[0]
            p_q.clear()
            _add_text_with_fraction_formatting(p_q, q_text, font_family, question_font_pt)
            p_q.space_after = Pt(int(8 * spacing_mult))
            p_q.space_before = Pt(int(4 * spacing_mult))
            # append options or answer lines
            if q['type'] == 'mcq' and opts:
                for opt in opts:
                    p = tfq.add_paragraph()
                    _add_text_with_fraction_formatting(p, opt, font_family, option_font_pt)
                    p.level = 0
                    p.space_before = Pt(int(4 * spacing_mult))
                    p.space_after = Pt(int(4 * spacing_mult))
            elif q['type'] == 'fill':
                p = tfq.add_paragraph()
                _add_text_with_fraction_formatting(
                    p,
                    existing_answer if existing_answer else 'Answer: __________________________',
                    font_family,
                    answer_font_pt,
                )
            else:
                p = tfq.add_paragraph()
                _add_text_with_fraction_formatting(
                    p,
                    existing_answer if existing_answer else 'Answer:',
                    font_family,
                    answer_font_pt,
                )
        else:
            # Content too long - split into chunks across multiple slides
            q_chunks = split_text_to_fit(q_text, question_font_pt, usable_width_in, usable_height_in)
            if not q_chunks:
                q_chunks = ['']

            for ci, chunk in enumerate(q_chunks):
                if ci > 0:
                    slide = prs.slides.add_slide(blank_layout)
                chunk_h_in = estimate_height_inches(chunk, question_font_pt, usable_width_in)
                question_height = Inches(max(0.6, chunk_h_in))
                txq = slide.shapes.add_textbox(place_left, place_top, prs.slide_width - left - right - extra_right, question_height)
                tfq = txq.text_frame
                # ensure wrapping and small internal margins so text wraps inside box
                try:
                    tfq.word_wrap = True
                except Exception:
                    pass
                try:
                    tfq.margin_left = Pt(6)
                    tfq.margin_right = Pt(6)
                    tfq.margin_top = Pt(6)
                    tfq.margin_bottom = Pt(6)
                except Exception:
                    pass
                tfq.clear()
                p_q = tfq.paragraphs[0]
                p_q.clear()
                _add_text_with_fraction_formatting(p_q, chunk, font_family, question_font_pt)
                p_q.space_after = Pt(int(8 * spacing_mult))
                p_q.space_before = Pt(int(4 * spacing_mult))

                if ci == len(q_chunks) - 1:
                    if q['type'] == 'mcq' and opts:
                        for opt in opts:
                            p = tfq.add_paragraph()
                            _add_text_with_fraction_formatting(p, opt, font_family, option_font_pt)
                            p.level = 0
                            p.space_before = Pt(int(4 * spacing_mult))
                            p.space_after = Pt(int(4 * spacing_mult))
                    elif q['type'] == 'fill':
                        p = tfq.add_paragraph()
                        _add_text_with_fraction_formatting(
                            p,
                            existing_answer if existing_answer else 'Answer: __________________________',
                            font_family,
                            answer_font_pt,
                        )
                    else:
                        p = tfq.add_paragraph()
                        _add_text_with_fraction_formatting(
                            p,
                            existing_answer if existing_answer else 'Answer:',
                            font_family,
                            answer_font_pt,
                        )

    # VERIFICATION STEP: Scan all slides and fix any font inconsistencies
    verify_and_fix_pptx(prs, font_family, question_font_pt, option_font_pt, answer_font_pt)
    
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')


@app.route('/convert', methods=['POST'])
def convert():
    f = request.files.get('file')
    if not f:
        return 'No file uploaded', 400
    filename = f.filename or ''
    lower = filename.lower()
    temp_paths = []
    try:
        # Prefer a Word copy/paste pipeline for DOCX on Windows to preserve fractions exactly.
        if lower.endswith('.docx') or f.mimetype == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            docx_path = save_upload_temp(f, suffix='.docx')
            temp_paths.append(docx_path)

            # Parse parameters early (COM path uses them)
            font_family = request.form.get('fontFamily', 'Calibri')
            try:
                font_size = int(request.form.get('fontSize', '36'))
            except (ValueError, TypeError):
                font_size = 36

            pptx_io = _create_pptx_by_copy_paste_from_docx(docx_path, font_family, font_size)
            if pptx_io is not None:
                base = os.path.splitext(filename)[0] if filename else ''
                safe_base = secure_filename(base) or 'converted'
                download_name = safe_base + '.pptx'
                return send_file(
                    pptx_io,
                    as_attachment=True,
                    download_name=download_name,
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )

            # Fallback: if COM isn't available, continue with the normal parser.
            with open(docx_path, 'rb') as docxf:
                title, questions = parse_docx(docxf)

        if lower.endswith('.pdf') or f.mimetype == 'application/pdf':
            title, questions = parse_pdf(f.stream)
        elif lower.endswith('.doc') or f.mimetype in ('application/msword', 'application/vnd.ms-word'):
            # save, convert, then parse
            path = save_upload_temp(f, suffix='.doc')
            temp_paths.append(path)
            try:
                converted = convert_doc_to_docx(path)
            except Exception as e:
                return f'Error converting .doc: {e}', 500
            with open(converted, 'rb') as convf:
                title, questions = parse_docx(convf)
        else:
            # treat as docx
            title, questions = parse_docx(f)
    finally:
        for p in temp_paths:
            try:
                os.remove(p)
            except Exception:
                pass
    
    # Parse form parameters with proper type conversion
    font_family = request.form.get('fontFamily', 'Calibri')
    try:
        font_size = int(request.form.get('fontSize', '36'))
    except (ValueError, TypeError):
        font_size = 36
    try:
        line_spacing = float(request.form.get('lineSpacing', '1.25'))
    except (ValueError, TypeError):
        line_spacing = 1.25
    
    pptx_io = create_pptx(
        title, 
        questions,
        font_family=font_family,
        font_size=font_size,
        line_spacing=line_spacing
    )
    # Use uploaded file's base name for download (sanitized), preserve .pptx
    base = os.path.splitext(filename)[0] if filename else ''
    safe_base = secure_filename(base) or 'converted'
    download_name = safe_base + '.pptx'
    return send_file(pptx_io, as_attachment=True, download_name=download_name, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    uploads_dir = os.path.join(os.path.dirname(__file__), 'uploads')
    return send_from_directory(uploads_dir, filename)


@app.route('/preview', methods=['POST'])
def preview():
    f = request.files.get('file')
    if not f:
        return jsonify({'error': 'no file'}), 400
    filename = secure_filename(f.filename or str(uuid.uuid4()))
    lower = filename.lower()
    uploads_dir = os.path.join(os.path.dirname(__file__), 'uploads')
    os.makedirs(uploads_dir, exist_ok=True)
    # Support .doc by converting to .docx first
    if lower.endswith('.pdf') or f.mimetype == 'application/pdf':
        save_name = f"{uuid.uuid4().hex}_{filename}"
        path = os.path.join(uploads_dir, save_name)
        f.save(path)
        return jsonify({'type': 'pdf', 'url': f'/uploads/{save_name}'}), 200
    elif lower.endswith('.doc') or f.mimetype in ('application/msword', 'application/vnd.ms-word'):
        # save and convert
        try:
            tmp = save_upload_temp(f, suffix='.doc')
            converted = convert_doc_to_docx(tmp)
            with open(converted, 'rb') as convf:
                title, questions = parse_docx(convf)
        except Exception as e:
            return jsonify({'error': 'could not convert .doc', 'msg': str(e)}), 400
        finally:
            try:
                os.remove(tmp)
            except Exception:
                pass

        snippet = ''
        if title:
            snippet += title + '\n\n'
        for q in questions[:5]:
            snippet += q.get('question', '')[:800] + '\n\n'
        return jsonify({'type': 'docx', 'text': snippet}), 200
    else:
        # try docx
        try:
            title, questions = parse_docx(f)
            snippet = ''
            if title:
                snippet += title + '\n\n'
            # include first few questions as preview
            for q in questions[:5]:
                snippet += q.get('question','')[:800] + '\n\n'
            return jsonify({'type': 'docx', 'text': snippet}), 200
        except Exception as e:
            return jsonify({'error': 'could not parse file', 'msg': str(e)}), 400


if __name__ == '__main__':
    # Default to stable dev-server settings; allow overrides via environment vars.
    host = os.environ.get('APP_HOST', '127.0.0.1')
    try:
        port = int(os.environ.get('APP_PORT', '5000'))
    except ValueError:
        port = 5000
    debug = os.environ.get('APP_DEBUG', '0') == '1'
    app.run(host=host, port=port, debug=debug)
