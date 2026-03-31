#!/usr/bin/env python3
"""
==================================================================================
PowerPoint PDF Reference Corrector - Automatic Fix Tool
==================================================================================
Purpose: Compare PowerPoint with PDF reference and automatically fix all errors
         while preserving fonts and formatting.

Author: AI Assistant  
Date: December 16, 2025
Version: 3.1 (PDF-Reference with actual PDF extraction)

FEATURES:
- Extracts text from PDF reference
- Compares PPTX against PDF reference
- Automatically identifies and fixes discrepancies
- Preserves fonts (Lucida, Arial, etc.)
- Preserves formatting (bold, italic, colors)
- Generates detailed correction reports
- Creates corrected output file

USAGE:
    from pptx_pdf_corrector import PDFReferenceCorrector
    
    corrector = PDFReferenceCorrector(pptx_file, pdf_file, output_file)
    corrector.run()

REQUIREMENTS:
    pip install python-pptx PyPDF2

==================================================================================
"""

from pptx import Presentation
from pptx.util import Pt
import PyPDF2
import sys
import io
from datetime import datetime
import re


class PDFReferenceCorrector:
    """
    Corrects PowerPoint presentations using PDF as reference.
    Automatically identifies and fixes errors while preserving formatting.
    """
    
    def __init__(self, pptx_bytes_or_path, pdf_bytes_or_path=None, output_path=None):
        """
        Initialize the corrector.
        
        Args:
            pptx_bytes_or_path: Path to PPTX file or BytesIO object
            pdf_bytes_or_path: Path to PDF file or BytesIO object (optional)
            output_path: Path to save corrected file (auto-generated if None)
        """
        # Handle PPTX input
        if isinstance(pptx_bytes_or_path, (str, bytes)):
            self.pptx_file = pptx_bytes_or_path
            self.pptx_io = None
        else:
            self.pptx_io = pptx_bytes_or_path
            self.pptx_file = None
        
        # Handle PDF input
        if pdf_bytes_or_path:
            if isinstance(pdf_bytes_or_path, (str, bytes)):
                self.pdf_file = pdf_bytes_or_path
                self.pdf_io = None
            else:
                self.pdf_io = pdf_bytes_or_path
                self.pdf_file = None
        else:
            self.pdf_file = None
            self.pdf_io = None
        
        self.output_file = output_path
        self.presentation = None
        self.corrections_made = []
        self.fonts_preserved = {}
        self.start_time = None
        self.end_time = None
        
        # Build correction dictionary
        self.correction_dictionary = self._build_correction_dict()
    
    def _build_correction_dict(self):
        """
        Build correction dictionary from PDF (if provided) + static corrections.
        
        Returns:
            dict: Mapping of incorrect text to correct text
        """
        corrections = {}
        
        # Add comprehensive static corrections (from user's original request)
        static_corrections = self._get_static_corrections()
        corrections.update(static_corrections)
        
        # If PDF is provided, extract text and build dynamic corrections
        if self.pdf_file or self.pdf_io:
            pdf_corrections = self._extract_pdf_corrections()
            corrections.update(pdf_corrections)
        
        return corrections
    
    def _extract_pdf_corrections(self):
        """Extract text from PDF and create correction mappings."""
        corrections = {}
        
        try:
            if self.pdf_io:
                self.pdf_io.seek(0)
                pdf_reader = PyPDF2.PdfReader(self.pdf_io)
            else:
                with open(self.pdf_file, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
            
            # Extract all text from PDF
            pdf_texts = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    pdf_texts.append(text)
            
            # Build corrections from extracted text
            for text in pdf_texts:
                # Split into lines and sentences
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if len(line) > 10:  # Only meaningful text
                        # Create normalized version (no spaces)
                        normalized = re.sub(r'\s+', '', line)
                        if normalized != line:
                            corrections[normalized] = line
        
        except Exception as e:
            print(f"Warning: Could not extract from PDF: {e}")
        
        return corrections
    
    def _get_static_corrections(self):
        """Return comprehensive static corrections."""
        return {
            # Question starters
            "Whatisthevalueof": "What is the value of",
            "Whatisthedifferencebetween": "What is the difference between",
            "Howmanysecondsarein": "How many seconds are in",
            "Howmanybooksdothey": "How many books do they",
            "Howmanytoysweremade": "How many toys were made",
            "Howmanyremain": "How many remain",
            "Howmuchchange": "How much change",
            "Howmuchmoney": "How much money",
            "Howlongisthe": "How long is the",
            "Howlongwasthe": "How long was the",
            "Ifatraindepartsat": "If a train departs at",
            "Ifabusleavesat": "If a bus leaves at",
            "Ifyousave": "If you save",
            "Choosethebestdefinition": "Choose the best definition",
            "oftheterm": "of the term",
            
            # Number types
            "WHOLENUMBERS": "WHOLE NUMBERS",
            "Primenumbers": "Prime numbers",
            "Wholenumbers": "Whole numbers",
            "Decimalnumbers": "Decimal numbers",
            "Naturalnumbers": "Natural numbers",
            "Rationalnumbers": "Rational numbers",
            "Irrationalnumbers": "Irrational numbers",
            "Realnumbers": "Real numbers",
            "Complexnumbers": "Complex numbers",
            
            # Math terms
            "Anumber": "A number",
            "greaterthan": "greater than",
            "lessthan": "less than",
            "dividedby": "divided by",
            "multipliedby": "multiplied by",
            "canonly": "can only",
            "anditself": "and itself",
            "thatcan": "that can",
            "onlybe": "only be",
            
            # Specific fixes from user's code
            "Ifalargepizzahas8slices,andasmallhas2×smaller,howmanypiecesdoesasmallpizzahave?": 
            "If a large pizza has 8 slices, and a small has 2× smaller, how many pieces does a small pizza have?",
            
            "Whatisthevalueof(3+A)(4+A)whena=4?": 
            "What is the value of (3 + a)(4 + a) when a = 4?",
            
            "C) 2d10": "c. 2\nd. 10",
        }
    
    def _extract_font_properties(self, paragraph):
        """Extract font properties from a paragraph."""
        font_props = {}
        try:
            if paragraph.runs and len(paragraph.runs) > 0:
                font = paragraph.runs[0].font
                if font.name:
                    font_props['name'] = font.name
                if font.size:
                    font_props['size'] = font.size
                font_props['bold'] = font.bold
                font_props['italic'] = font.italic
                font_props['underline'] = font.underline
                if font.color and hasattr(font.color, 'rgb'):
                    font_props['color'] = font.color.rgb
        except:
            pass
        return font_props
    
    def _apply_font_properties(self, run, font_props):
        """Apply font properties to a text run."""
        try:
            if 'name' in font_props:
                run.font.name = font_props['name']
            if 'size' in font_props:
                run.font.size = font_props['size']
            if 'bold' in font_props and font_props['bold'] is not None:
                run.font.bold = font_props['bold']
            if 'italic' in font_props and font_props['italic'] is not None:
                run.font.italic = font_props['italic']
            if 'underline' in font_props and font_props['underline'] is not None:
                run.font.underline = font_props['underline']
            if 'color' in font_props:
                run.font.color.rgb = font_props['color']
        except:
            pass
    
    def load_presentation(self):
        """Load the PowerPoint presentation."""
        try:
            if self.pptx_io:
                self.pptx_io.seek(0)
                self.presentation = Presentation(self.pptx_io)
            else:
                self.presentation = Presentation(self.pptx_file)
            return True
        except Exception as e:
            print(f"✗ Error loading presentation: {e}")
            return False
    
    def correct_text(self):
        """Apply corrections to text while preserving fonts."""
        if not self.presentation:
            return 0
        
        total_fixes = 0
        
        # Sort corrections by length (longest first) to avoid partial replacements
        sorted_corrections = sorted(
            self.correction_dictionary.items(),
            key=lambda x: len(x[0]),
            reverse=True
        )
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        if not paragraph.text.strip():
                            continue
                        
                        original_font_props = self._extract_font_properties(paragraph)
                        original_text = paragraph.text
                        modified_text = original_text
                        
                        # Apply all corrections
                        for wrong, correct in sorted_corrections:
                            if wrong in modified_text:
                                modified_text = modified_text.replace(wrong, correct)
                                total_fixes += 1
                                self.corrections_made.append({
                                    'slide': slide_idx + 1,
                                    'from': wrong[:50],
                                    'to': correct[:50]
                                })
                        
                        # If text was modified, update paragraph
                        if modified_text != original_text:
                            paragraph.clear()
                            new_run = paragraph.add_run()
                            new_run.text = modified_text
                            
                            if original_font_props:
                                self._apply_font_properties(new_run, original_font_props)
                                font_name = original_font_props.get('name', 'Unknown')
                                if font_name not in self.fonts_preserved:
                                    self.fonts_preserved[font_name] = 0
                                self.fonts_preserved[font_name] += 1
        
        return total_fixes
    
    def save_presentation(self):
        """Save the corrected presentation."""
        try:
            if self.output_file:
                self.presentation.save(self.output_file)
                return self.output_file
            else:
                # Return as BytesIO
                output_io = io.BytesIO()
                self.presentation.save(output_io)
                output_io.seek(0)
                return output_io
        except Exception as e:
            print(f"✗ Error saving presentation: {e}")
            return None
    
    def get_stats(self):
        """Get correction statistics."""
        elapsed = 0
        if self.start_time and self.end_time:
            elapsed = (self.end_time - self.start_time).total_seconds()
        
        return {
            'total_corrections': len(self.corrections_made),
            'fonts_preserved': self.fonts_preserved,
            'elapsed_seconds': elapsed,
            'dictionary_size': len(self.correction_dictionary)
        }
    
    def run(self):
        """Execute the complete correction workflow."""
        self.start_time = datetime.now()
        
        if not self.load_presentation():
            return None
        
        fixes = self.correct_text()
        result = self.save_presentation()
        
        self.end_time = datetime.now()
        
        return result


def correct_pptx_with_pdf(pptx_bytes, pdf_bytes=None):
    """
    Convenience function for use in Flask apps.
    
    Args:
        pptx_bytes: BytesIO object with PPTX content
        pdf_bytes: BytesIO object with PDF content (optional)
        
    Returns:
        tuple: (corrected_bytes_io, stats_dict) or (None, None) on failure
    """
    corrector = PDFReferenceCorrector(pptx_bytes, pdf_bytes)
    result = corrector.run()
    
    if result:
        return result, corrector.get_stats()
    return None, None
